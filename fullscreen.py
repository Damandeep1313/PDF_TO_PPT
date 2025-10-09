import os
import re
import json
import time
import hashlib
import random
import httpx
import requests
from io import BytesIO
import uuid
import flask
from importlib.metadata import version
from functools import wraps



from dotenv import load_dotenv
from concurrent.futures import ThreadPoolExecutor, as_completed
from tqdm import tqdm

from flask import Flask, request, jsonify
from datetime import datetime
from openai import OpenAI
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import traceback
import os
from io import BytesIO
from gpt_image_generator import ImageGenerator


import cloudinary
import cloudinary.uploader
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

# ----------- CONFIG -----------
load_dotenv()
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
CLOUDINARY_CLOUD_NAME = os.getenv("CLOUDINARY_CLOUD_NAME")
CLOUDINARY_API_KEY = os.getenv("CLOUDINARY_API_KEY")
CLOUDINARY_API_SECRET = os.getenv("CLOUDINARY_API_SECRET")

if not OPENAI_API_KEY:
    raise ValueError("OPENAI_API_KEY not set in environment variables")
if not (CLOUDINARY_CLOUD_NAME and CLOUDINARY_API_KEY and CLOUDINARY_API_SECRET):
    raise ValueError("Cloudinary credentials missing in .env")

client = OpenAI(
    api_key=OPENAI_API_KEY,
    # Explicitly prevent proxy interference:
    http_client=httpx.Client(trust_env=False)  # ← ADD THIS
)
cloudinary.config(
    cloud_name=CLOUDINARY_CLOUD_NAME,
    api_key=CLOUDINARY_API_KEY,
    api_secret=CLOUDINARY_API_SECRET,
    secure=True
)

IMG_SIZE = "1024x1024"
MAX_WORKERS = 4
IMAGE_CACHE_DIR = "img_cache"
os.makedirs(IMAGE_CACHE_DIR, exist_ok=True)

PROFESSIONAL_PALETTES = [
    {
        "name": "Corporate Blue",
        "primary": "#1E3A8A", "secondary": "#3B82F6", "accent": "#60A5FA",
        "text": "#FFFFFF", "text_dark": "#1F2937", "gradient_start": "#1E3A8A", "gradient_end": "#3B82F6"
    },
    {
        "name": "Elegant Gray",
        "primary": "#4B5563", "secondary": "#6B7280", "accent": "#9CA3AF",
        "text": "#F9FAFB", "text_dark": "#111827", "gradient_start": "#4B5563", "gradient_end": "#9CA3AF"
    },
    {
        "name": "Forest Green",
        "primary": "#065F46", "secondary": "#10B981", "accent": "#34D399",
        "text": "#ECFDF5", "text_dark": "#064E3B", "gradient_start": "#065F46", "gradient_end": "#10B981"
    },
    {
        "name": "Sunset Orange",
        "primary": "#C2410C", "secondary": "#F97316", "accent": "#FB923C",
        "text": "#FFF7ED", "text_dark": "#7C2D12", "gradient_start": "#C2410C", "gradient_end": "#F97316"
    },
    {
        "name": "Royal Purple",
        "primary": "#5B21B6", "secondary": "#8B5CF6", "accent": "#A78BFA",
        "text": "#F3E8FF", "text_dark": "#3B0764", "gradient_start": "#5B21B6", "gradient_end": "#8B5CF6"
    },
    {
        "name": "Ocean Teal",
        "primary": "#0F766E", "secondary": "#14B8A6", "accent": "#2DD4BF",
        "text": "#E0F2FE", "text_dark": "#134E4A", "gradient_start": "#0F766E", "gradient_end": "#14B8A6"
    },
    {
        "name": "Warm Sand",
        "primary": "#92400E", "secondary": "#D97706", "accent": "#F59E0B",
        "text": "#FFF8E1", "text_dark": "#78350F", "gradient_start": "#92400E", "gradient_end": "#D97706"
    },
    {
        "name": "Modern Slate",
        "primary": "#1E293B", "secondary": "#334155", "accent": "#64748B",
        "text": "#F1F5F9", "text_dark": "#0F172A", "gradient_start": "#1E293B", "gradient_end": "#334155"
    },
    {
        "name": "Deep Crimson",
        "primary": "#7F1D1D", "secondary": "#B91C1C", "accent": "#EF4444",
        "text": "#FEF2F2", "text_dark": "#4B0505", "gradient_start": "#7F1D1D", "gradient_end": "#B91C1C"
    },
    {
        "name": "Cool Indigo",
        "primary": "#4338CA", "secondary": "#6366F1", "accent": "#818CF8",
        "text": "#EEF2FF", "text_dark": "#312E81", "gradient_start": "#4338CA", "gradient_end": "#6366F1"
    },
    {
        "name": "Fresh Lime",
        "primary": "#365314", "secondary": "#4ADE80", "accent": "#A7F3D0",
        "text": "#F0FDF4", "text_dark": "#1C2F0E", "gradient_start": "#365314", "gradient_end": "#4ADE80"
    },
    {
        "name": "Midnight Black",
        "primary": "#111827", "secondary": "#374151", "accent": "#6B7280",
        "text": "#F9FAFB", "text_dark": "#000000", "gradient_start": "#111827", "gradient_end": "#374151"
    },
    {
        "name": "Soft Coral",
        "primary": "#BE123C", "secondary": "#F43F5E", "accent": "#FCA5A5",
        "text": "#FFF1F2", "text_dark": "#831843", "gradient_start": "#BE123C", "gradient_end": "#F43F5E"
    },
    {
        "name": "Steel Blue",
        "primary": "#1E40AF", "secondary": "#3B82F6", "accent": "#60A5FA",
        "text": "#E0E7FF", "text_dark": "#1E3A8A", "gradient_start": "#1E40AF", "gradient_end": "#3B82F6"
    },
    {
        "name": "Bright Cyan",
        "primary": "#0E7490", "secondary": "#22D3EE", "accent": "#67E8F9",
        "text": "#ECFEFF", "text_dark": "#164E63", "gradient_start": "#0E7490", "gradient_end": "#22D3EE"
    },
    {
        "name": "Goldenrod",
        "primary": "#B45309", "secondary": "#FBBF24", "accent": "#FCD34D",
        "text": "#FFFBEB", "text_dark": "#78350F", "gradient_start": "#B45309", "gradient_end": "#FBBF24"
    },
    {
        "name": "Classic Navy",
        "primary": "#0C4A6E", "secondary": "#2563EB", "accent": "#60A5FA",
        "text": "#EFF6FF", "text_dark": "#1E3A8A", "gradient_start": "#0C4A6E", "gradient_end": "#2563EB"
    },
    {
        "name": "Rich Burgundy",
        "primary": "#6B0218", "secondary": "#9F1239", "accent": "#DC2626",
        "text": "#FEE2E2", "text_dark": "#4B0109", "gradient_start": "#6B0218", "gradient_end": "#9F1239"
    },
    {
        "name": "Vibrant Orange",
        "primary": "#C2410C", "secondary": "#F97316", "accent": "#FDBA74",
        "text": "#FFF7ED", "text_dark": "#7C2D12", "gradient_start": "#C2410C", "gradient_end": "#F97316"
    },
    {
        "name": "Dusty Rose",
        "primary": "#881337", "secondary": "#BE185D", "accent": "#F472B6",
        "text": "#FFF1F2", "text_dark": "#4B0630", "gradient_start": "#881337", "gradient_end": "#BE185D"
    }
    
]

# ----------- HELPERS -----------





def optimize_layout(slide, content_type):  
    # Clears default placeholders, sets smart margins  
    for shape in slide.shapes:  
        if shape.is_placeholder:  
            shape.element.getparent().remove(shape.element)  

    # Content-aware positioning (title vs. bullet vs. image)  
    return {  
        'title': {'top': 0.5, 'left': 1.0},  
        'bullet': {'top': 2.0, 'left': 1.5},  
        'image': {'top': 1.5, 'left': 5.0}  
    }[content_type]  


def set_auto_text_color(shape, bg_color):
    """Automatically sets black or white text based on background brightness"""
    try:
        rgb = hex_to_rgb(bg_color)
        brightness = (rgb[0]*299 + rgb[1]*587 + rgb[2]*114) / 1000
        shape.text_frame.paragraphs[0].font.color.rgb = (
            RGBColor(0, 0, 0) if brightness > 128 
            else RGBColor(255, 255, 255))
    except Exception as e:
        print(f"⚠️ Text color error: {e}")
        # Fallback to black text
        shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)


def add_auto_cropped_image(slide, img_path, x, y, w, h):  
    img = slide.shapes.add_picture(img_path, x, y, w, h)  
    img.crop_left = img.crop_right = 0.1  # 10% auto-crop  
    img.crop_top = img.crop_bottom = 0.1  



def get_safe_font():
    """Returns available fonts in priority order"""
    for font in ["Calibri", "Arial", "Helvetica", "Segoe UI"]:
        if font in Presentation().font_manager:
            return font
    return "Calibri"




def add_random_design_element(slide, theme):
    """Adds random design elements with contrasting colors to slides"""
    palette = PROFESSIONAL_PALETTES[theme.get("palette_index", random.randint(0, len(PROFESSIONAL_PALETTES)-1))]
    
    # Available shapes (MSO_SHAPE enum values)
    shapes = [
        MSO_SHAPE.RECTANGLE,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        MSO_SHAPE.OVAL,
        MSO_SHAPE.DIAMOND,
        MSO_SHAPE.CHEVRON,
        MSO_SHAPE.PENTAGON,
        MSO_SHAPE.PLAQUE
    ]
    
    # Choose random properties
    shape_type = random.choice(shapes)
    rotation = random.randint(-45, 45)
    width = Inches(random.uniform(0.5, 3))
    height = Inches(random.uniform(0.1, 0.5))
    x_pos = Inches(random.uniform(-1, 10))
    y_pos = Inches(random.uniform(0, 7))
    
    # Choose contrasting color
    color_choices = [palette["accent"], palette["secondary"]]
    if random.random() > 0.7:
        color_choices.append(palette["text"])
    fill_color = random.choice(color_choices)
    
    # Create the shape
    shape = slide.shapes.add_shape(
        shape_type,
        x_pos, y_pos,
        width, height
    )
    
    # Apply styling
    shape.rotation = rotation
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*hex_to_rgb(fill_color))
    shape.fill.transparency = random.uniform(0.2, 0.6)
    shape.line.fill.background()  # No border
    
    # Correct shadow implementation
    if random.random() > 0.5:
        shadow = shape.shadow
        shadow.inherit = False
        shadow.visible = True
        shadow.blur_radius = Pt(5)
        shadow.offset_x = Pt(2)
        shadow.offset_y = Pt(2)
        # Set color through foreground color
        shadow.fore_color.rgb = RGBColor(0, 0, 0)
        shadow.transparency = 0.5




def add_premium_design_elements(slide, theme):
    """Adds professional design elements with precise placement"""
    palette = PROFESSIONAL_PALETTES[theme.get("palette_index", random.randint(0, len(PROFESSIONAL_PALETTES)-1))]
    
    # Main diagonal accent strip (perfectly aligned from corner to corner)
    diagonal = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(-2), Inches(0),  # Starts above top-left corner
        Inches(12), Inches(0.8)  # Long enough to cross entire slide
    )
    diagonal.rotation = -20  # Gentle angle
    diagonal.fill.solid()
    diagonal.fill.fore_color.rgb = RGBColor(*hex_to_rgb(palette["accent"]))
    diagonal.line.fill.background()
    diagonal.fill.transparency = 0.15
    
    # Secondary elements (precisely placed)
    elements = [
        # (x, y, width, height, rotation, color_key, transparency)
        (8.5, 0.5, 1.5, 0.3, 15, "secondary", 0.2),  # Top-right
        (0.5, 5.5, 2.0, 0.4, -15, "primary", 0.25),   # Bottom-left
        (6.0, 6.0, 1.0, 0.2, 0, "accent", 0.3)       # Bottom-center
    ]
    
    for x, y, w, h, rot, color_key, trans in elements:
        elem = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(w), Inches(h)
        )
        elem.rotation = rot
        elem.fill.solid()
        elem.fill.fore_color.rgb = RGBColor(*hex_to_rgb(palette[color_key]))
        elem.line.fill.background()
        elem.fill.transparency = trans





def add_safe_shadow(shape):
    """Bulletproof shadow implementation"""
    try:
        if not hasattr(shape, 'shadow'):
            return
            
        shadow = shape.shadow
        shadow.inherit = False
        shadow.visible = True
        shadow.blur_radius = Pt(4)
        shadow.offset_x = Pt(1)
        shadow.offset_y = Pt(1)
        
        # Universal color setting approach
        try:
            shadow._element.get_or_add_effectClr().srgbClr.val = '1F1F1F'  # Hex for dark gray
        except:
            try:
                if hasattr(shadow, 'color'):
                    shadow.color.rgb = RGBColor(31, 31, 31)
                elif hasattr(shadow.fill, 'fore_color'):
                    shadow.fill.solid()
                    shadow.fill.fore_color.rgb = RGBColor(31, 31, 31)
            except:
                pass
                
        shadow.transparency = 0.6
    except Exception as e:
        print(f"⚠️ Shadow formatting skipped: {str(e)}")



def clean_code_fence(s: str) -> str:
    s = s.strip()
    s = re.sub(r"^```(?:json)?\n", "", s)
    s = re.sub(r"\n```$", "", s)
    return s.strip()

def hex_to_rgb(hexstr):
    h = hexstr.lstrip("#")
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))

def add_professional_gradient(slide, start_color, end_color, direction="vertical"):
    try:
        bg_fill = slide.background.fill
        bg_fill.gradient()
        bg_fill.gradient_stops[0].color.rgb = RGBColor(*hex_to_rgb(start_color))
        bg_fill.gradient_stops[1].color.rgb = RGBColor(*hex_to_rgb(end_color))
        if direction == "diagonal":
            bg_fill.gradient_angle = 45
        elif direction == "horizontal":
            bg_fill.gradient_angle = 0
        else:
            bg_fill.gradient_angle = 90
    except Exception:
        try:
            bg_fill = slide.background.fill
            bg_fill.solid()
            bg_fill.fore_color.rgb = RGBColor(*hex_to_rgb(start_color))
        except:
            pass

# def create_professional_text_box(slide, x, y, width, height, text, theme,
#                                font_size=18, font_name="Calibri", alignment=PP_ALIGN.LEFT,
#                                bold=False, text_color_key="text"):
#     text_box = slide.shapes.add_textbox(x, y, width, height)
#     tf = text_box.text_frame
#     tf.clear()
#     tf.word_wrap = True
#     tf.auto_size = None
#     tf.vertical_anchor = MSO_ANCHOR.TOP
#     tf.margin_bottom = Pt(8)
#     tf.margin_top = Pt(8)
#     tf.margin_left = Pt(16)
#     tf.margin_right = Pt(16)
    
#     p = tf.paragraphs[0]
#     p.text = text
#     p.font.size = Pt(font_size)
#     p.font.name = font_name
#     p.font.bold = bold
#     p.alignment = alignment
#     p.space_after = Pt(15)
#     p.line_spacing = 1.3
    
#     # Color handling (updated)
#     if text_color_key == "auto":
#         # Get background color (assuming gradient start as bg)
#         bg_color = theme.get("gradient_start", "#FFFFFF")  
#         p.font.color.rgb = get_contrast_color(bg_color)
#     else:
#         p.font.color.rgb = RGBColor(*hex_to_rgb(theme.get(text_color_key, "#000000")))
    
#     return text_box

# Add this helper function
def get_contrast_color(bg_color):
    """Returns black or white depending on background brightness"""
    rgb = hex_to_rgb(bg_color)
    brightness = (rgb[0]*299 + rgb[1]*587 + rgb[2]*114) / 1000
    return RGBColor(0, 0, 0) if brightness > 128 else RGBColor(255, 255, 255)




def create_professional_shape(slide, shape_type, x, y, width, height, fill_color, transparency=0):
    shape = slide.shapes.add_shape(shape_type, x, y, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*hex_to_rgb(fill_color))
    if transparency > 0:
        shape.fill.transparency = transparency
    shape.line.fill.background()
    return shape

# ----------- CLASSES -----------

class EnhancedSlidePlanner:
    def __init__(self, client):
        self.client = client

    def plan_slides(self, doc_text, target_slide_count):
        prompt = f"""
You are a professional presentation designer creating a corporate-level presentation.

STRUCTURE REQUIREMENTS:
1. Title slide (handled separately)
2. Table of Contents slide (handled separately)  
3. Exactly {target_slide_count} content slides with substantial information
4. Images will be added to every 2nd content slide automatically

CONTENT REQUIREMENTS:
- Each slide must have comprehensive, professional content
- Titles: Clear, descriptive, professional (30-80 characters)
- Content: 3-5 substantial bullet points per slide (40-120 characters each)
- Professional tone throughout
- Each slide should cover a distinct topic/section

INPUT DOCUMENT:
{doc_text}

OUTPUT: JSON with this structure:
{{
  "presentation_meta": {{
    "title": "Professional presentation title",
    "subtitle": "Descriptive subtitle explaining the content",
    "total_content_slides": {target_slide_count},
    "estimated_duration": "{target_slide_count * 2}-{target_slide_count * 3} minutes"
  }},
  "theme": {{
    "name": "Professional Theme Name",
    "style": "corporate",
    "palette_index": 0,
    "mood": "professional"
  }},
  "table_of_contents": [
    {{
      "section_number": 1,
      "section_title": "First main section title",
      "slides": [1, 2]
    }},
    {{
      "section_number": 2,
      "section_title": "Second main section title", 
      "slides": [3, 4]
    }}
  ],
  "content_slides": [
    {{
      "slide_number": 1,
      "section": "Introduction",
      "title": "Professional slide title that clearly describes the content",
      "content_points": [
        "First comprehensive point with detailed explanation and context",
        "Second substantial point providing valuable insights and information", 
        "Third detailed point with specific examples and actionable content",
        "Fourth comprehensive point that adds significant value to understanding"
      ],
      "slide_type": "text_heavy",
      "has_image": false,
      "image_concept": "Professional concept for image if needed"
    }},
    {{
      "slide_number": 2,
      "section": "Introduction", 
      "title": "Second slide title with clear professional focus",
      "content_points": [
        "Detailed first point with comprehensive explanation and examples",
        "Substantial second point providing deep insights and practical value",
        "Third comprehensive point with specific data and actionable recommendations",
        "Fourth detailed point that enhances understanding and provides clarity"
      ],
      "slide_type": "image_slide",
      "has_image": true,
      "image_concept": "Professional, clean image concept that supports the slide content and maintains corporate aesthetic"
    }}
  ]
}}

Create exactly {target_slide_count} content slides with professional, substantial content.
Ensure every 2nd slide (slides 2, 4, 6, 8, etc.) has has_image: true.
Return only valid JSON.
"""
        resp = self.client.chat.completions.create(
            model="gpt-4o",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.7
        )
        raw = resp.choices[0].message.content
        raw = clean_code_fence(raw)
        try:
            data = json.loads(raw)
            slides = data.get("content_slides", [])
            for i, slide in enumerate(slides):
                slide_num = slide.get("slide_number", i + 1)
                slide["has_image"] = (slide_num % 2 == 0)
                slide["slide_type"] = "image_slide" if slide["has_image"] else "text_heavy"
                if slide["has_image"] and not slide.get("image_concept"):
                    slide["image_concept"] = f"Professional illustration representing {slide.get('title', 'slide content')}, clean corporate style, modern design"
            return data.get("presentation_meta", {}), data.get("theme", {}), data.get("table_of_contents", []), slides
        except Exception as e:
            print(f"JSON parse error: {e}")
            return {}, {}, [], []

class ProfessionalImageGenerator:
    def __init__(self, client, max_workers=MAX_WORKERS, cache_dir=IMAGE_CACHE_DIR):
        self.client = client
        self.max_workers = max_workers
        self.cache_dir = cache_dir

    def _prompt_to_filename(self, prompt):
        h = hashlib.sha256(prompt.encode("utf-8")).hexdigest()
        return os.path.join(self.cache_dir, f"{h}.png")

    def enhance_professional_prompt(self, prompt):
        professional_style = """
professional corporate design, clean minimal aesthetic, high-end business presentation style,
sophisticated color scheme, premium quality, modern flat design, no text overlays,
suitable for executive presentation, clean background, professional photography style,
corporate branding appropriate, business-focused imagery
"""
        return f"{prompt}, {professional_style}"

    def generate_image(self, prompt):
        filename = self._prompt_to_filename(prompt)
        if os.path.exists(filename):
            print(f"📁 Using cached image for: {prompt[:50]}...")
            return filename
        enhanced_prompt = self.enhance_professional_prompt(prompt)
        for attempt in range(3):
            try:
                print(f"🎨 Generating professional image: {prompt[:50]}...")
                resp = self.client.images.generate(
                    model="dall-e-3",
                    prompt=enhanced_prompt,
                    size=IMG_SIZE,
                    quality="hd"
                )
                if hasattr(resp, 'data') and resp.data:
                    entry = resp.data[0]
                    if hasattr(entry, 'url') and entry.url:
                        img_bytes = requests.get(entry.url).content
                        with open(filename, "wb") as f:
                            f.write(img_bytes)
                        return filename
                time.sleep(2)
            except Exception as e:
                print(f"⚠️ Image generation error attempt {attempt+1}: {e}")
                time.sleep(2 ** attempt)
        return None

    def generate_images_for_slides(self, slides):
        print("🖼️ Generating professional images for designated slides...")
        image_paths = {}
        slides_needing_images = [s for s in slides if s.get("has_image")]
        with ThreadPoolExecutor(max_workers=self.max_workers) as executor:
            futures = {}
            for slide in slides_needing_images:
                slide_num = slide.get("slide_number")
                prompt = slide.get("image_concept", f"Professional illustration for {slide.get('title', 'slide')}")
                futures[executor.submit(self.generate_image, prompt)] = slide_num
            for fut in tqdm(as_completed(futures), total=len(futures), desc="🎨 Creating professional visuals"):
                slide_num = futures[fut]
                result = fut.result()
                if result:
                    image_paths[slide_num] = result
        print(f"✅ Generated {len(image_paths)} professional images")
        return image_paths

class ProfessionalPPTBuilder:
    def __init__(self):
        # Initialize default styling parameters
        self.default_title_size = Pt(44)
        self.default_subtitle_size = Pt(24)
        self.default_font = 'Calibri'


    def create_professional_text_box(self, slide, x, y, width, height, text, theme,
                              font_size=18, font_name="Calibri", alignment=PP_ALIGN.LEFT,
                              bold=False, text_color_key="text"):
        text_box = slide.shapes.add_textbox(x, y, width, height)
        tf = text_box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.auto_size = None
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_bottom = Pt(8)
        tf.margin_top = Pt(8)
        tf.margin_left = Pt(16)
        tf.margin_right = Pt(16)
        
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.name = font_name
        p.font.bold = bold
        p.alignment = alignment
        p.space_after = Pt(15)
        p.line_spacing = 1.3
    
    # Handle color assignment safely
        try:
            if isinstance(theme, dict):  # If theme is a palette dictionary
                if text_color_key == "auto":
                    bg_color = theme.get("gradient_start", "#FFFFFF")
                    p.font.color.rgb = self.get_contrast_color(bg_color)
                else:
                    color_hex = theme.get(text_color_key, "#000000")
                    p.font.color.rgb = RGBColor(*hex_to_rgb(color_hex))
            else:  # If theme is already an RGBColor or similar
                p.font.color.rgb = theme
        except Exception as e:
            print(f"⚠️ Color assignment error: {e}")
            p.font.color.rgb = RGBColor(0, 0, 0)  # Fallback to black
        
        return text_box


    def get_contrast_color(self, bg_color):
        """Returns black or white depending on background brightness"""
        rgb = hex_to_rgb(bg_color)
        brightness = (rgb[0]*299 + rgb[1]*587 + rgb[2]*114) / 1000
        return RGBColor(0, 0, 0) if brightness > 128 else RGBColor(255, 255, 255)

    def create_professional_shape(self, slide, shape_type, x, y, width, height, fill_color, transparency=0):
        shape = slide.shapes.add_shape(shape_type, x, y, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*hex_to_rgb(fill_color))
        if transparency > 0:
            shape.fill.transparency = transparency
        shape.line.fill.background()
        return shape


    def _create_boxed_text_layout(self, slide, slide_data, palette):
        """Creates text content in colored boxes (without title)"""
        points = slide_data.get("content_points", ["Point 1", "Point 2", "Point 3"])
        if not points:  # Ensure we always have content
            points = ["Key content not specified"]
        
        # Box colors from palette
        box_colors = [palette["primary"], palette["secondary"], palette["accent"]]
        
        # Create boxes in 2x2 grid (starting lower since title exists)
        for i, point in enumerate(points[:4]):  # Max 4 boxes
            row = i // 2
            col = i % 2
            left = Inches(0.8 + col * 4.5)
            top = Inches(2.5 + row * 2)  # Adjusted to start below title
            
            # Create rounded box
            box = self.create_professional_shape(
                slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                left, top, Inches(4), Inches(1.8),
                box_colors[i % len(box_colors)],
                transparency=0.2
            )
            
            # Add text with automatic contrast
            bg_color = box_colors[i % len(box_colors)]
            text_color = self.get_contrast_color(bg_color)
            self.create_professional_text_box(
                slide, left + Inches(0.3), top + Inches(0.3),
                Inches(3.4), Inches(1.2), point,
                text_color, font_size=16
            )







    def _apply_corporate_design(self, prs, theme):
        """Internal method to set master slide styles"""
        try:
            # Get color palette
            palette = PROFESSIONAL_PALETTES[theme.get("palette_index", 0)]
            
            # Set master background
            background = prs.slide_master.background
            background.fill.solid()
            background.fill.fore_color.rgb = RGBColor(*hex_to_rgb(palette["gradient_start"]))
            
            # Set default text styles through placeholders
            for layout in prs.slide_master.slide_layouts:
                try:
                    placeholders = layout.placeholders
                    
                    # Title placeholder (usually index 0)
                    if len(placeholders) > 0:
                        title = placeholders[0]
                        if hasattr(title, 'text_frame'):
                            title.text_frame.paragraphs[0].font.name = self.default_font
                            title.text_frame.paragraphs[0].font.size = Pt(36)
                            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(*hex_to_rgb(palette["text"]))
                    
                    # Body placeholder (usually index 1)
                    if len(placeholders) > 1:
                        body = placeholders[1]
                        if hasattr(body, 'text_frame'):
                            for paragraph in body.text_frame.paragraphs:
                                paragraph.font.name = self.default_font
                                paragraph.font.size = Pt(18)
                                paragraph.font.color.rgb = RGBColor(*hex_to_rgb(palette["text_dark"]))
                except Exception as e:
                    print(f"⚠️ Layout styling error: {e}")
                    continue
                    
        except Exception as e:
            print(f"⚠️ Design template error: {e}")


    
    def create_title_slide(self, prs, presentation_meta, theme):
        """Creates a professional title slide with perfect spacing and responsive design"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        palette = PROFESSIONAL_PALETTES[theme.get("palette_index", random.randint(0, len(PROFESSIONAL_PALETTES)-1))]
        
        # 1. Background Design
        add_professional_gradient(slide, palette["gradient_start"], palette["gradient_end"], "diagonal")
        
        # Add premium design elements (behind text)
        self.add_premium_title_elements(slide, palette)

        # 2. Dynamic Title Configuration
        title_text = presentation_meta.get("title", "Professional Presentation")
        title_font_size = self.calculate_optimal_font_size(title_text, max_chars=50)
        
        self.create_professional_text_box(
            slide, 
            Inches(0.5), Inches(1.8),  # Position (left, top)
            Inches(9), Inches(1.8),     # Dimensions (width, height)
            title_text,
            palette,
            font_size=title_font_size,
            alignment=PP_ALIGN.CENTER,
            bold=True,
            text_color_key="text"
        )

        # 3. Smart Subtitle Handling
        subtitle_text = presentation_meta.get("subtitle", "Comprehensive Analysis")
        subtitle_lines = self.split_text_to_lines(subtitle_text, max_line_length=60)
        
        self.create_professional_text_box(
            slide,
            Inches(0.5), Inches(4.0),  # Positioned lower than title
            Inches(9), Inches(len(subtitle_lines) * 0.6),  # Dynamic height
            '\n'.join(subtitle_lines),
            palette,
            font_size=20,
            alignment=PP_ALIGN.CENTER,
            text_color_key="text"
        )

        # 4. Decorative Elements
        self.add_title_slide_decoration(slide, palette)


    def add_premium_title_elements(self, slide, palette):
        """Adds design elements specifically for title slide"""
        # Diagonal accent strip
        diagonal = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(-2), Inches(1),
            Inches(12), Inches(0.6)
        )
        diagonal.rotation = -20
        diagonal.fill.solid()
        diagonal.fill.fore_color.rgb = RGBColor(*hex_to_rgb(palette["accent"]))
        diagonal.fill.transparency = 0.2
        diagonal.line.fill.background()

        # Corner elements
        for x, y, w, h, rot in [(8.5, 0.5, 1.5, 0.3, 15), (0.5, 6.5, 2.0, 0.4, -15)]:
            elem = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y), Inches(w), Inches(h)
            )
            elem.rotation = rot
            elem.fill.solid()
            elem.fill.fore_color.rgb = RGBColor(*hex_to_rgb(palette["secondary"]))
            elem.fill.transparency = 0.25
            elem.line.fill.background()


    def add_title_slide_decoration(self, slide, palette):
        """Adds decorative elements to title slide"""
        # Thin center line
        create_professional_shape(
            slide,
            random.choice([MSO_SHAPE.ROUNDED_RECTANGLE, MSO_SHAPE.OVAL]),
            Inches(2.5), Inches(5.5),
            Inches(5), Inches(0.05),
            palette["primary"]
        )
        
        # Slide number placeholder (even on title slide for consistency)
        create_professional_shape(
            slide,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(8.8), Inches(6.8),
            Inches(0.8), Inches(0.4),
            palette["accent"]
        )


    def calculate_optimal_font_size(self, text, max_chars):
        """Dynamically adjusts font size based on text length"""
        length = len(text)
        if length <= 30: return 44
        if length <= 50: return 38
        if length <= 70: return 32
        return 28


    def split_text_to_lines(self, text, max_line_length=60):
        """Smart text splitting for subtitles"""
        words = text.split()
        lines = []
        current_line = []
        
        for word in words:
            if len(' '.join(current_line + [word])) <= max_line_length:
                current_line.append(word)
            else:
                lines.append(' '.join(current_line))
                current_line = [word]
        
        if current_line:
            lines.append(' '.join(current_line))
            
        return lines                




    def create_toc_slide(self, prs, toc_data, theme):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        palette = PROFESSIONAL_PALETTES[theme.get("palette_index", random.randint(0, len(PROFESSIONAL_PALETTES)-1))]
        add_professional_gradient(slide, palette["gradient_start"], palette["gradient_end"])
        self.create_professional_text_box(
            slide, Inches(1), Inches(0.8), Inches(8), Inches(1),
            "Table of Contents", palette, font_size=36, alignment=PP_ALIGN.CENTER,
            bold=True, text_color_key="text"
        )
        create_professional_shape(
            slide, MSO_SHAPE.RECTANGLE,
            Inches(3), Inches(1.9), Inches(4), Inches(0.05),
            palette["accent"]
        )
        y_start = 2.5
        for i, section in enumerate(toc_data):
            section_num = section.get("section_number", i + 1)
            section_title = section.get("section_title", f"Section {section_num}")
            slides_range = section.get("slides", [])
            slide_range_text = f"Slides {slides_range[0]}-{slides_range[-1]}" if slides_range else f"Slide {section_num}"
            create_professional_shape(
                slide, MSO_SHAPE.OVAL,
                Inches(1.5), Inches(y_start + i * 0.8), Inches(0.6), Inches(0.6),
                palette["accent"]
            )
            self.create_professional_text_box(
                slide, Inches(1.5), Inches(y_start + i * 0.8), Inches(0.6), Inches(0.6),
                str(section_num), {"text": palette["text"]}, font_size=20,
                alignment=PP_ALIGN.CENTER, bold=True, text_color_key="text"
            )
            self.create_professional_text_box(
                slide, Inches(2.5), Inches(y_start + i * 0.8), Inches(5.5), Inches(0.6),
                section_title, palette, font_size=18, bold=True, text_color_key="text"
            )
            self.create_professional_text_box(
                slide, Inches(8.2), Inches(y_start + i * 0.8), Inches(1.5), Inches(0.6),
                slide_range_text, palette, font_size=14, alignment=PP_ALIGN.RIGHT,
                text_color_key="text"
            )


    def create_content_slide(self, slide, slide_data, image_path, theme, slide_number):
        palette = PROFESSIONAL_PALETTES[theme.get("palette_index", random.randint(0, len(PROFESSIONAL_PALETTES)-1))]
        add_professional_gradient(slide, palette["gradient_start"], palette["gradient_end"])
        
        # Add premium design elements (before content)
        add_premium_design_elements(slide, theme)
        
        # Slide number indicator
        create_professional_shape(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(9), Inches(0.2), Inches(0.8), Inches(0.4),
            palette["accent"]
        )
        self.create_professional_text_box(
            slide, Inches(9), Inches(0.2), Inches(0.8), Inches(0.4),
            str(slide_number), {"text": palette["text"]}, font_size=16,
            alignment=PP_ALIGN.CENTER, bold=True, text_color_key="text"
        )

        # Add title (only once, here in the main method)
        self.create_professional_text_box(
            slide, Inches(9), Inches(0.2), Inches(0.8), Inches(0.4),
            str(slide_number), {"text": palette["text"]}, font_size=16,
            alignment=PP_ALIGN.CENTER, bold=True, text_color_key="text"
        )

        # Add title with tighter line spacing
        title_box = self.create_professional_text_box(
            slide, Inches(0.8), Inches(0.5), Inches(8), Inches(1),
            slide_data.get("title", "Untitled Slide"), palette,
            font_size=28, bold=True, text_color_key="text"
        )
        
        # Add underline below title
        # create_professional_shape(
        #     slide, MSO_SHAPE.RECTANGLE,
        #     Inches(0.8), Inches(1.4), Inches(3), Inches(0.05),
        #     palette["accent"]
        # )

        # Handle content
        content_points = slide_data.get("content_points", ["Key points not specified"])
        if not content_points:  # Double-check empty list
            content_points = ["Important content goes here"]

        # Image handling
        has_image = (slide_data.get("has_image", False) and 
                    image_path and 
                    os.path.exists(image_path))
        
        if has_image:
            try:
                # Add image (right side)
                img = slide.shapes.add_picture(
                    image_path, 
                    Inches(5.5), Inches(1.8),  # x, y (below title)
                    Inches(3.5), Inches(3.5)    # width, height
                )
                
                # Add bullet points (left side)
                y_pos = 2.0  # Start position below title
                for point in content_points[:4]:  # Max 4 points
                    # Bullet marker
                    create_professional_shape(
                        slide, MSO_SHAPE.OVAL,
                        Inches(0.8), Inches(y_pos), Inches(0.15), Inches(0.15),
                        palette["accent"]
                    )
                    # Bullet text
                    self.create_professional_text_box(
                        slide, Inches(1.1), Inches(y_pos-0.08), 
                        Inches(4), Inches(0.5), point,
                        palette, font_size=14, text_color_key="text"
                    )
                    y_pos += 0.7  # Spacing between points
                    
            except Exception as e:
                print(f"⚠️ Image load failed, using text layout: {str(e)}")
                self._create_text_slide_layout(slide, slide_data, palette)
        else:
            # Text-only slide layouts
            if random.random() < 0.3:
                self._create_boxed_text_layout(slide, slide_data, palette)
            else:
                self._create_text_slide_layout(slide, slide_data, palette)

    def _create_boxed_text_layout(self, slide, slide_data, palette):
        """Creates content in colored boxes with corrected dimensions"""
        points = slide_data.get("content_points", ["Point 1", "Point 2", "Point 3", "Point 4"])
        if not points:
             points = ["Key content not specified"]

        # Box colors from palette
        box_colors = [palette["primary"], palette["secondary"], palette["accent"], palette["secondary"]]

        # Define grid properties that fit within the slide
        num_cols = 2
        
        # New, smaller dimensions and corrected positioning
        box_width = Inches(4.2)
        box_height = Inches(1.5)
        start_left = Inches(0.7)
        start_top = Inches(1.8) # Lowered to give title more space
        horz_gap = Inches(0.2)
        vert_gap = Inches(0.2)

        # Create boxes in a 2x2 grid
        for i, point in enumerate(points[:4]):  # Max 4 boxes
            row = i // num_cols
            col = i % num_cols
            
            # Calculate position based on new properties
            left = start_left + col * (box_width + horz_gap)
            top = start_top + row * (box_height + vert_gap)

            # Create rounded box with new dimensions
            box = self.create_professional_shape(
                slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                left, top, box_width, box_height, # Use new dimensions
                box_colors[i % len(box_colors)],
                transparency=0.15
            )

            # Add text with automatic contrast, slightly padded
            text_color = self.get_contrast_color(box_colors[i % len(box_colors)])
            self.create_professional_text_box(
                slide, left + Inches(0.2), top + Inches(0.2),
                box_width - Inches(0.4), box_height - Inches(0.4), # Text box fits inside shape
                point,
                text_color, 
                font_size=14, # Slightly smaller font for better fit
                alignment=PP_ALIGN.CENTER
            )

    def _create_text_slide_layout(self, slide, slide_data, palette):
        """Creates standard bullet points (without title)"""
        # 1. Ensure we always have content
        content_points = slide_data.get("content_points", ["Key points not specified"])
        if not content_points or not any(content_points):  # Handle empty lists
            content_points = ["Important content goes here"]
        
        # 2. Bullet points (starting below where title would be)
        y_start = 2.2  # Start position accounts for title space
        for i, point in enumerate(content_points[:6]):  # Max 6 points
            # Skip empty points
            if not point.strip():
                continue
                
            # Bullet marker
            create_professional_shape(
                slide, MSO_SHAPE.OVAL,
                Inches(1), Inches(y_start + i * 0.7),
                Inches(0.15), Inches(0.15),
                palette["accent"]
            )
            
            # Bullet text
            self.create_professional_text_box(
                slide, 
                Inches(1.4), Inches(y_start + i * 0.7 - 0.1),
                Inches(7), Inches(0.6),
                point, 
                palette, 
                font_size=16, 
                text_color_key="text"
            )
            
            # Prevent overflow
            if y_start + (i * 0.7) > 6.5:
                break




    def _create_image_slide_layout(self, slide, slide_data, palette):
        """Fallback layout when image fails to load"""
        y_start = 2.2
        content_points = slide_data.get("content_points", [])
        
        for i, point in enumerate(content_points[:4]):  # Max 4 points
            create_professional_shape(
                slide, MSO_SHAPE.OVAL,
                Inches(1), Inches(y_start + i * 0.7), Inches(0.12), Inches(0.12),
                palette["accent"]
            )
            self.create_professional_text_box(
                slide, Inches(1.3), Inches(y_start + i * 0.7 - 0.1), 
                Inches(7), Inches(0.5), point,
                palette, font_size=14, text_color_key="text"
            )




    def build(self, presentation_meta, theme, toc_data, slides, image_paths):
            prs = Presentation()
            prs.slide_width = Inches(10.0) 
            prs.slide_height = Inches(5.625)
            self._apply_corporate_design(prs, theme)
            
            self.create_title_slide(prs, presentation_meta, theme)
            self.create_toc_slide(prs, toc_data, theme)
            
            for slide_data in slides:
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                slide_num = slide_data.get("slide_number", 1)
                image_path = image_paths.get(slide_num)
                presentation_slide_num = slide_num + 2
                print(f"Creating slide {presentation_slide_num}: {slide_data.get('title', 'Untitled')}")
                self.create_content_slide(slide, slide_data, image_path, theme, presentation_slide_num)
                
            ppt_bytes_io = BytesIO()
            prs.save(ppt_bytes_io)
            ppt_bytes_io.seek(0)
            return ppt_bytes_io
     


def api_key_required(f):
    @wraps(f)
    def decorated(*args, **kwargs):
        api_key = request.headers.get('X-API-KEY')
        valid_keys = os.getenv('API_KEYS', '').split(',')
        
        if not api_key or api_key not in valid_keys:
            return jsonify({
                "status": "error",
                "error": "Invalid or missing API key"
            }), 403
        return f(*args, **kwargs)
    return decorated
    
# ----------- FLASK APP -----------



app = Flask(__name__)


def generate_presentation(slide_count, summary_text):
    try:
        # Initialize components
        planner = EnhancedSlidePlanner(client)
        image_gen = ImageGenerator(api_key=OPENAI_API_KEY, max_workers=10)  # Updated
        builder = ProfessionalPPTBuilder()

        # 1. Plan slides
        presentation_meta, theme, toc_data, slides = planner.plan_slides(summary_text, slide_count)
        if not slides or len(slides) < slide_count:
            raise ValueError(f"Failed to generate adequate slides (requested: {slide_count}, got: {len(slides) if slides else 0})")

        # 2. Set theme and generate images
        theme["palette_index"] = random.randint(0, len(PROFESSIONAL_PALETTES) - 1)
        
        # Get image prompts from slides that need images
        slides_needing_images = [s for s in slides if s.get("has_image")]
        image_prompts = {
            s["slide_number"]: s.get("image_concept", f"Image for slide {s['slide_number']}")
            for s in slides_needing_images
        }
        
        # Generate images concurrently (returns {prompt: path})
        generated_images = image_gen.generate_images(list(image_prompts.values()))
        
        # Map back to slide numbers {slide_num: image_path}
        image_paths = {
            slide_num: generated_images[prompt]
            for slide_num, prompt in image_prompts.items()
            if prompt in generated_images
        }
                
        # 3. Build presentation
        ppt_bytes_io = builder.build(presentation_meta, theme, toc_data, slides, image_paths)
        
        # Validate presentation
        if ppt_bytes_io.getbuffer().nbytes < 1024:
            raise ValueError("Generated presentation is too small (likely empty)")

        # 4. Create local backup
        local_path = f"presentation_{int(time.time())}.pptx"
        with open(local_path, "wb") as f:
            f.write(ppt_bytes_io.getvalue())
        print(f"✓ Local backup saved: {os.path.abspath(local_path)}")

        # 5. Upload to Cloudinary
        try:
            upload_result = cloudinary.uploader.upload(
                local_path,
                resource_type="raw",
                public_id=f"ppt/presentation_{int(time.time())}",
                overwrite=True
            )
            os.remove(local_path)
            return upload_result["secure_url"]
        except Exception as upload_error:
            print(f"⚠️ Cloudinary upload failed: {upload_error}")
            return f"file://{os.path.abspath(local_path)}"

    except Exception as e:
        print(f"❌ Generation failed: {e}\n{traceback.format_exc()}")
        raise RuntimeError(f"Presentation generation failed: {e}") from e







@app.route('/')
def home():
    return "Service is live ✅"
@app.route("/generate-ppt", methods=["POST"])
@api_key_required  # Add this line exactly here

def generate_ppt_endpoint():
    """Generate PowerPoint from JSON input with comprehensive error handling"""
    try:
        # 1. Validate Input
        if not request.is_json:
            return jsonify({
                "error": "Content-Type must be application/json",
                "status": "invalid_request"
            }), 400

        data = request.get_json()
        if not data:
            return jsonify({
                "error": "Empty JSON body",
                "status": "invalid_request"
            }), 400

        # 2. Extract and Validate Parameters
        try:
            slide_count = int(data.get("slide_count", 0))
            if slide_count < 1 or slide_count > 20:  # Reasonable limit
                raise ValueError
        except (TypeError, ValueError):
            return jsonify({
                "error": "slide_count must be an integer between 1-20",
                "status": "invalid_parameter"
            }), 400

        summary = data.get("summary", "").strip()
        if len(summary) < 20:  # Minimum length check
            return jsonify({
                "error": "summary must be at least 20 characters",
                "status": "invalid_parameter"
            }), 400

        # 3. Generate Presentation
        ppt_url = generate_presentation(slide_count, summary)
        
        # Determine if URL is local or cloud
        is_local = ppt_url.startswith("file://")
        
        return jsonify({
            "status": "success",
            "url": ppt_url,
            "source": "local" if is_local else "cloudinary",
            "timestamp": datetime.now().isoformat(),
            "slide_count": slide_count
        })

    except Exception as e:
        app.logger.error(f"PPT Generation Failed: {str(e)}\n{traceback.format_exc()}")
        return jsonify({
            "status": "error",
            "error": "Presentation generation failed",
            "details": str(e),
            "trace_id": str(uuid.uuid4())  # For support tracking
        }), 500

if __name__ == "__main__":
    # Configure production-ready settings
    app.config['JSON_SORT_KEYS'] = False
    app.config['MAX_CONTENT_LENGTH'] = 10 * 1024 * 1024  # 10MB limit
    
    print(f"🚀 Starting PPT Generator (Flask {version('flask')})")
    app.run(
        host='0.0.0.0', 
        port=5001,
        debug=False,  # Never use debug=True in production
        threaded=True
    )
